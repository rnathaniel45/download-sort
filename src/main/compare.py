from sentence_transformers import SentenceTransformer, util
from PIL import Image
import pytesseract
import sys
import easyocr
import torch
import json
import os
from PyPDF2 import PdfReader
from docx import Document
model = SentenceTransformer('all-mpnet-base-v2')

def ocr_image(image_path, folders):
    reader = easyocr.Reader(['en'],verbose=False)  # English model
    result = reader.readtext(image_path)
    latex_lines = [entry[1] for entry in result]  # get only the recognized text
    latex_code = " ".join(latex_lines)
    files = [latex_code]

    # Embed folder names and files
    folder_embeddings = model.encode(list(map(lambda f : f["desc"], folders)), convert_to_tensor=True, show_progress_bar=False)
    file_embeddings = model.encode(files, convert_to_tensor=True, show_progress_bar=False)

    # Compute cosine similarity between each file and all folder embeddings
    cosine_scores = util.cos_sim(file_embeddings, folder_embeddings)

    # Assign each file to the folder with highest similarity
    best_match_index = torch.argmax(cosine_scores)

    # Get the corresponding folder name
    best_match_folder = folders[best_match_index]["path"]

    return best_match_folder

def powerpoint_ocr(image_path, folders): 
    from pptx import Presentation
    prs = Presentation(image_path)
    text = ""

    for slide in prs.slides:
                for shape in slide.shapes:
                    # Text directly in shapes
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    # Text inside tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                text += cell.text + "\n"
    folder_embeddings = model.encode(list(map(lambda f : f["desc"], folders)), convert_to_tensor=True, show_progress_bar=False)
    file_embeddings = model.encode(text, convert_to_tensor=True, show_progress_bar=False)

    cosine_scores = util.cos_sim(file_embeddings, folder_embeddings)

    best_match_index = torch.argmax(cosine_scores)

    return folders[best_match_index]["path"]

def pdf_ocr(image_path, folders):
    with open(image_path, "rb") as file:
        reader = PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        folder_embeddings = model.encode(list(map(lambda f : f["desc"], folders)), convert_to_tensor=True, show_progress_bar=False)
        file_embeddings = model.encode(text, convert_to_tensor=True, show_progress_bar=False)

        cosine_scores = util.cos_sim(file_embeddings, folder_embeddings)

        best_match_index = torch.argmax(cosine_scores)

        return folders[best_match_index]["path"]
def docx_ocr(image_path, folders):
    doc = Document(image_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    folder_embeddings = model.encode(list(map(lambda f : f["desc"], folders)), convert_to_tensor=True, show_progress_bar=False)
    file_embeddings = model.encode(text, convert_to_tensor=True, show_progress_bar=False)

    cosine_scores = util.cos_sim(file_embeddings, folder_embeddings)

    best_match_index = torch.argmax(cosine_scores)

    return folders[best_match_index]["path"]

# # Load the pretrained MPNet model


# Define your folder categories (labels)


# Example file names to classify

if __name__ == "__main__":
    func = sys.argv[1]
    folders = json.loads(sys.argv[2])
    if(os.path.splitext(func)[1] == ".pptx"):
        print(powerpoint_ocr(func, folders))
    elif(os.path.splitext(func)[1] == ".pdf"):
        print(pdf_ocr(func, folders))
    elif(os.path.splitext(func)[1] == ".docx"):
        print(docx_ocr(func, folders))
    else:
        print(ocr_image(func, folders))
